"""Render strict JSON-like requests to DOCX, XLSX, or PPTX bytes.

The contract is intentionally deterministic and network-free: an LLM or host
application authors a JSON payload, while this module only validates and
renders it. Strings written to spreadsheets are treated as literal text by
default, preventing formula injection from AI-authored content.
"""

from __future__ import annotations

import json
import math
import re
from collections.abc import Mapping, Sequence
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation

from .docx_image import DocxImageContractError, add_docx_inline_png


class OfficeDocumentError(ValueError):
    """Raised when an Office rendering request violates the public contract."""


@dataclass(frozen=True, slots=True)
class RenderedOfficeDocument:
    """In-memory Office artifact and its transport metadata."""

    format: str
    extension: str
    content_type: str
    data: bytes


_FORMATS = {
    "docx": (
        ".docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ),
    "xlsx": (
        ".xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ),
    "pptx": (
        ".pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ),
}
_INVALID_SHEET_NAME = re.compile(r"[\\/*?:\[\]]")
_FORMULA_PREFIXES = ("=", "+", "-", "@")
_MAX_DOCX_RICH_RUNS = 4096
_DOCX_PARAGRAPH_ALIGNMENTS = {
    "left": WD_ALIGN_PARAGRAPH.LEFT,
    "center": WD_ALIGN_PARAGRAPH.CENTER,
    "right": WD_ALIGN_PARAGRAPH.RIGHT,
    "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
}
_MISSING = object()


def load_schema() -> dict[str, Any]:
    """Load the bundled JSON Schema used for structured LLM output."""

    text = Path(__file__).with_name("schema.json").read_text(encoding="utf-8")
    return json.loads(text)


def render_office_document(payload: Mapping[str, Any]) -> RenderedOfficeDocument:
    """Validate and render one Office document request entirely in memory."""

    request = _mapping(payload, "payload")
    format_name = _string(_require(request, "format", "payload"), "format")
    if format_name not in _FORMATS:
        raise OfficeDocumentError(f"unsupported format: {format_name!r}")

    if format_name == "docx":
        data = _render_docx(request)
    elif format_name == "xlsx":
        data = _render_xlsx(request)
    else:
        data = _render_pptx(request)

    extension, content_type = _FORMATS[format_name]
    return RenderedOfficeDocument(format_name, extension, content_type, data)


def _render_docx(request: Mapping[str, Any]) -> bytes:
    """Render a validated DOCX request to an in-memory OOXML package."""

    _reject_unknown(request, {"format", "title", "author", "subject", "blocks"}, "payload")
    document = Document()
    title = _optional_string(request.get("title", _MISSING), "title")
    author = _optional_string(request.get("author", _MISSING), "author")
    subject = _optional_string(request.get("subject", _MISSING), "subject")
    if title:
        document.core_properties.title = title
        document.add_heading(title, level=0)
    if author:
        document.core_properties.author = author
    if subject:
        document.core_properties.subject = subject

    blocks = _array(_require(request, "blocks", "payload"), "blocks")
    for index, raw_block in enumerate(blocks):
        path = f"blocks[{index}]"
        block = _mapping(raw_block, path)
        block_type = _string(_require(block, "type", path), f"{path}.type")
        if block_type == "heading":
            _reject_unknown(block, {"type", "text", "level", "alignment"}, path)
            level = _integer(
                _require(block, "level", path),
                f"{path}.level",
                minimum=1,
                maximum=9,
            )
            heading = document.add_heading(
                _string(_require(block, "text", path), f"{path}.text"),
                level=level,
            )
            _apply_docx_paragraph_alignment(heading, block, path)
        elif block_type == "paragraph":
            _reject_unknown(block, {"type", "text", "alignment"}, path)
            paragraph = document.add_paragraph(
                _string(
                    _require(block, "text", path),
                    f"{path}.text",
                    allow_empty=True,
                )
            )
            _apply_docx_paragraph_alignment(paragraph, block, path)
        elif block_type == "rich_paragraph":
            _reject_unknown(block, {"type", "runs", "alignment"}, path)
            _add_docx_rich_paragraph(document, block, path)
        elif block_type == "bullet_list":
            _reject_unknown(block, {"type", "items", "ordered"}, path)
            items = _array(_require(block, "items", path), f"{path}.items")
            ordered = _boolean(block.get("ordered", False), f"{path}.ordered")
            style = "List Number" if ordered else "List Bullet"
            for item_index, item in enumerate(items):
                document.add_paragraph(
                    _string(item, f"{path}.items[{item_index}]", allow_empty=True),
                    style=style,
                )
        elif block_type == "table":
            _reject_unknown(block, {"type", "headers", "rows"}, path)
            _add_docx_table(document, block, path)
        elif block_type == "image":
            _reject_unknown(block, {"type", "source", "alt_text", "width_px"}, path)
            try:
                add_docx_inline_png(document, block, path)
            except DocxImageContractError as exc:
                raise OfficeDocumentError(str(exc)) from None
        elif block_type == "page_break":
            _reject_unknown(block, {"type"}, path)
            document.add_page_break()
        else:
            raise OfficeDocumentError(f"{path}.type is unsupported: {block_type!r}")

    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _apply_docx_paragraph_alignment(
    paragraph: Any, block: Mapping[str, Any], path: str
) -> None:
    """Apply one explicit bounded Word paragraph alignment when present."""

    if "alignment" not in block:
        return
    alignment_path = f"{path}.alignment"
    alignment = _string(block["alignment"], alignment_path)
    try:
        paragraph.alignment = _DOCX_PARAGRAPH_ALIGNMENTS[alignment]
    except KeyError:
        allowed = ", ".join(_DOCX_PARAGRAPH_ALIGNMENTS)
        raise OfficeDocumentError(
            f"{alignment_path} must be one of: {allowed}"
        ) from None


def _add_docx_rich_paragraph(
    document: Document, block: Mapping[str, Any], path: str
) -> None:
    """Append one bounded paragraph of explicitly formatted deterministic runs."""

    runs = _array(_require(block, "runs", path), f"{path}.runs")
    if not runs:
        raise OfficeDocumentError(f"{path}.runs must contain at least one run")
    if len(runs) > _MAX_DOCX_RICH_RUNS:
        raise OfficeDocumentError(
            f"{path}.runs must contain at most {_MAX_DOCX_RICH_RUNS} runs"
        )

    paragraph = document.add_paragraph()
    _apply_docx_paragraph_alignment(paragraph, block, path)
    for run_index, raw_run in enumerate(runs):
        run_path = f"{path}.runs[{run_index}]"
        run_spec = _mapping(raw_run, run_path)
        _reject_unknown(run_spec, {"text", "bold", "italic", "underline"}, run_path)
        text = _string(
            _require(run_spec, "text", run_path),
            f"{run_path}.text",
            allow_empty=True,
        )
        if text == "":
            raise OfficeDocumentError(f"{run_path}.text must not be empty")
        run = paragraph.add_run(text)
        if "bold" in run_spec:
            run.bold = _boolean(run_spec["bold"], f"{run_path}.bold")
        if "italic" in run_spec:
            run.italic = _boolean(run_spec["italic"], f"{run_path}.italic")
        if "underline" in run_spec:
            run.underline = _boolean(run_spec["underline"], f"{run_path}.underline")


def _add_docx_table(document: Document, block: Mapping[str, Any], path: str) -> None:
    """Append a validated, rectangular table block to a Word document."""

    headers = _array(block.get("headers", []), f"{path}.headers")
    rows = _array(_require(block, "rows", path), f"{path}.rows")
    normalized_headers = [
        _scalar(value, f"{path}.headers[{index}]") for index, value in enumerate(headers)
    ]
    normalized_rows = _normalize_rows(rows, f"{path}.rows")
    width = len(normalized_headers) or (len(normalized_rows[0]) if normalized_rows else 0)
    if width == 0:
        raise OfficeDocumentError(f"{path} table must contain at least one column")
    if normalized_headers and any(len(row) != width for row in normalized_rows):
        raise OfficeDocumentError(f"{path} rows must have the same width as headers")
    if not normalized_headers and any(len(row) != width for row in normalized_rows):
        raise OfficeDocumentError(f"{path} rows must have the same width")

    table = document.add_table(
        rows=len(normalized_rows) + (1 if normalized_headers else 0),
        cols=width,
    )
    table.style = "Table Grid"
    offset = 0
    if normalized_headers:
        for column, value in enumerate(normalized_headers):
            table.cell(0, column).text = _display(value)
        offset = 1
    for row_index, row in enumerate(normalized_rows, start=offset):
        for column, value in enumerate(row):
            table.cell(row_index, column).text = _display(value)


def _render_xlsx(request: Mapping[str, Any]) -> bytes:
    """Render a validated workbook request to an in-memory OOXML package."""

    _reject_unknown(request, {"format", "title", "author", "sheets"}, "payload")
    sheets = _array(_require(request, "sheets", "payload"), "sheets")
    if not sheets:
        raise OfficeDocumentError("sheets must contain at least one worksheet")

    workbook = Workbook()
    workbook.remove(workbook.active)
    title = _optional_string(request.get("title", _MISSING), "title")
    author = _optional_string(request.get("author", _MISSING), "author")
    if title:
        workbook.properties.title = title
    if author:
        workbook.properties.creator = author

    seen_names: set[str] = set()
    for sheet_index, raw_sheet in enumerate(sheets):
        path = f"sheets[{sheet_index}]"
        sheet = _mapping(raw_sheet, path)
        _reject_unknown(
            sheet,
            {"name", "rows", "header_row", "freeze_panes", "auto_filter"},
            path,
        )
        name = _string(_require(sheet, "name", path), f"{path}.name")
        _validate_sheet_name(name, path, seen_names)
        seen_names.add(name.casefold())

        worksheet = workbook.create_sheet(name)
        rows = _normalize_rows(
            _array(_require(sheet, "rows", path), f"{path}.rows"),
            f"{path}.rows",
        )
        for row_index, row in enumerate(rows, start=1):
            for column_index, value in enumerate(row, start=1):
                cell = worksheet.cell(row=row_index, column=column_index)
                cell.value = value
                if isinstance(value, str) and value.lstrip().startswith(_FORMULA_PREFIXES):
                    cell.data_type = "s"

        header_row = _boolean(sheet.get("header_row", False), f"{path}.header_row")
        if header_row and rows:
            for cell in worksheet[1]:
                cell.font = Font(bold=True)
                cell.fill = PatternFill(fill_type="solid", fgColor="D9EAF7")

        freeze_panes = _optional_string(
            sheet.get("freeze_panes", _MISSING), f"{path}.freeze_panes"
        )
        if freeze_panes:
            worksheet.freeze_panes = freeze_panes

        auto_filter = _boolean(sheet.get("auto_filter", False), f"{path}.auto_filter")
        if auto_filter and rows and max((len(row) for row in rows), default=0):
            worksheet.auto_filter.ref = worksheet.dimensions

        _size_columns(worksheet, rows)

    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


def _validate_sheet_name(name: str, path: str, seen_names: set[str]) -> None:
    """Reject worksheet names that Excel cannot address unambiguously."""

    if len(name) > 31 or _INVALID_SHEET_NAME.search(name):
        raise OfficeDocumentError(f"{path}.name is invalid for Excel: {name!r}")
    if name.casefold() in seen_names:
        raise OfficeDocumentError("worksheet names must be unique (case-insensitive)")


def _size_columns(worksheet: Any, rows: list[list[Any]]) -> None:
    """Apply bounded content-derived widths to populated worksheet columns."""

    width = max((len(row) for row in rows), default=0)
    for column_index in range(1, width + 1):
        longest = max(
            (len(_display(row[column_index - 1])) for row in rows if len(row) >= column_index),
            default=0,
        )
        worksheet.column_dimensions[get_column_letter(column_index)].width = min(
            max(longest + 2, 10), 60
        )


def _render_pptx(request: Mapping[str, Any]) -> bytes:
    """Render a validated slide-deck request to an in-memory OOXML package."""

    _reject_unknown(request, {"format", "title", "author", "slides"}, "payload")
    slides = _array(_require(request, "slides", "payload"), "slides")
    if not slides:
        raise OfficeDocumentError("slides must contain at least one slide")

    presentation = Presentation()
    title = _optional_string(request.get("title", _MISSING), "title")
    author = _optional_string(request.get("author", _MISSING), "author")
    if title:
        presentation.core_properties.title = title
    if author:
        presentation.core_properties.author = author

    for slide_index, raw_slide in enumerate(slides):
        path = f"slides[{slide_index}]"
        slide_spec = _mapping(raw_slide, path)
        _reject_unknown(slide_spec, {"title", "subtitle", "bullets"}, path)
        slide_title = _string(
            _require(slide_spec, "title", path), f"{path}.title"
        )
        subtitle = _optional_string(
            slide_spec.get("subtitle", _MISSING), f"{path}.subtitle"
        )
        bullets = _array(slide_spec.get("bullets", []), f"{path}.bullets")
        if subtitle is not None and "bullets" in slide_spec:
            raise OfficeDocumentError(f"{path}.subtitle cannot be combined with bullets")

        if subtitle is not None:
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            slide.shapes.title.text = slide_title
            slide.placeholders[1].text = subtitle
            continue

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_title
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for bullet_index, bullet in enumerate(bullets):
            text, level = _normalize_bullet(bullet, f"{path}.bullets[{bullet_index}]")
            paragraph = text_frame.paragraphs[0] if bullet_index == 0 else text_frame.add_paragraph()
            paragraph.text = text
            paragraph.level = level

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _normalize_bullet(value: Any, path: str) -> tuple[str, int]:
    """Normalize one string-or-object bullet into text and nesting level."""

    if isinstance(value, str):
        return value, 0
    if isinstance(value, Mapping):
        _reject_unknown(value, {"text", "level"}, path)
        text = _string(
            _require(value, "text", path), f"{path}.text", allow_empty=True
        )
        level = _integer(value.get("level", 0), f"{path}.level", minimum=0, maximum=8)
        return text, level
    raise OfficeDocumentError(f"{path} must be a string or object")


def _require(mapping: Mapping[str, Any], key: str, path: str) -> Any:
    """Return a required mapping field or raise a path-qualified error."""

    if key not in mapping:
        raise OfficeDocumentError(f"{path}.{key} is required")
    return mapping[key]


def _reject_unknown(
    mapping: Mapping[str, Any], allowed: set[str], path: str
) -> None:
    """Reject undeclared fields so generated requests cannot be ambiguous."""

    unexpected = sorted(set(mapping) - allowed)
    if unexpected:
        label = "field" if len(unexpected) == 1 else "fields"
        raise OfficeDocumentError(
            f"{path} has unexpected {label}: {', '.join(unexpected)}"
        )


def _mapping(value: Any, path: str) -> Mapping[str, Any]:
    """Validate and return a string-keyed mapping value."""

    if not isinstance(value, Mapping):
        raise OfficeDocumentError(f"{path} must be an object")
    if any(not isinstance(key, str) for key in value):
        raise OfficeDocumentError(f"{path} object keys must be strings")
    return value


def _array(value: Any, path: str) -> list[Any]:
    """Validate and materialize a non-string sequence as a list."""

    if not isinstance(value, Sequence) or isinstance(value, (str, bytes, bytearray)):
        raise OfficeDocumentError(f"{path} must be an array")
    return list(value)


def _string(value: Any, path: str, *, allow_empty: bool = False) -> str:
    """Validate a string, optionally permitting empty or whitespace-only text."""

    if not isinstance(value, str):
        raise OfficeDocumentError(f"{path} must be a string")
    if not allow_empty and not value.strip():
        raise OfficeDocumentError(f"{path} must not be empty")
    return value


def _optional_string(value: Any, path: str) -> str | None:
    """Return an absent optional string as ``None`` and validate present values."""

    if value is _MISSING:
        return None
    return _string(value, path)


def _boolean(value: Any, path: str) -> bool:
    """Validate and return a strict JSON boolean."""

    if not isinstance(value, bool):
        raise OfficeDocumentError(f"{path} must be a boolean")
    return value


def _integer(value: Any, path: str, *, minimum: int, maximum: int) -> int:
    """Validate a non-boolean integer inside an inclusive range."""

    if not isinstance(value, int) or isinstance(value, bool):
        raise OfficeDocumentError(f"{path} must be an integer")
    if value < minimum or value > maximum:
        raise OfficeDocumentError(f"{path} must be between {minimum} and {maximum}")
    return value


def _scalar(value: Any, path: str) -> str | int | float | bool | None:
    """Validate a finite JSON scalar suitable for an Office cell or table."""

    if isinstance(value, float) and not math.isfinite(value):
        raise OfficeDocumentError(f"{path} must be a finite JSON number")
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    raise OfficeDocumentError(f"{path} must be a JSON scalar")


def _normalize_rows(rows: list[Any], path: str) -> list[list[Any]]:
    """Validate nested row arrays and normalize every cell to a JSON scalar."""

    normalized: list[list[Any]] = []
    for row_index, raw_row in enumerate(rows):
        row = _array(raw_row, f"{path}[{row_index}]")
        normalized.append(
            [_scalar(value, f"{path}[{row_index}][{column}]") for column, value in enumerate(row)]
        )
    return normalized


def _display(value: Any) -> str:
    """Convert an Office table value to display text without spelling ``None``."""

    return "" if value is None else str(value)
