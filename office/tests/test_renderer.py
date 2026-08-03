from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from inkspan_office import (
    OfficeDocumentError,
    load_schema,
    render_office_document,
    write_office_document,
)


def test_load_schema_exposes_all_supported_formats() -> None:
    schema = load_schema()

    assert schema["$schema"].endswith("2020-12/schema")
    assert schema["title"] == "Inkspan Office document request"
    assert {entry["properties"]["format"]["const"] for entry in schema["oneOf"]} == {
        "docx",
        "xlsx",
        "pptx",
    }
    assert schema["$defs"]["slide"]["not"]["required"] == ["subtitle", "bullets"]
    assert schema["$defs"]["nonEmptyString"]["pattern"] == r"\S"
    assert schema["oneOf"][0]["properties"]["title"]["$ref"] == (
        "#/$defs/nonEmptyString"
    )


def test_render_docx_supports_ai_authored_blocks_and_metadata() -> None:
    rendered = render_office_document(
        {
            "format": "docx",
            "title": "Quarterly brief",
            "author": "Inkspan Agent",
            "subject": "Board update",
            "blocks": [
                {"type": "heading", "level": 1, "text": "Executive summary"},
                {"type": "paragraph", "text": "Revenue grew while churn fell."},
                {
                    "type": "bullet_list",
                    "ordered": False,
                    "items": ["Ship the beta", "Measure adoption"],
                },
                {
                    "type": "bullet_list",
                    "ordered": True,
                    "items": ["Approve budget"],
                },
                {
                    "type": "table",
                    "headers": ["Metric", "Value"],
                    "rows": [["Revenue", 120], ["Churn", 0.03]],
                },
                {"type": "page_break"},
                {"type": "paragraph", "text": "Appendix"},
            ],
        }
    )

    assert rendered.format == "docx"
    assert rendered.extension == ".docx"
    assert rendered.content_type.endswith("wordprocessingml.document")
    assert rendered.data.startswith(b"PK")

    document = Document(BytesIO(rendered.data))
    assert document.core_properties.title == "Quarterly brief"
    assert document.core_properties.author == "Inkspan Agent"
    assert document.core_properties.subject == "Board update"
    assert [p.text for p in document.paragraphs][:5] == [
        "Quarterly brief",
        "Executive summary",
        "Revenue grew while churn fell.",
        "Ship the beta",
        "Measure adoption",
    ]
    assert document.paragraphs[3].style.name == "List Bullet"
    assert document.paragraphs[5].style.name == "List Number"
    assert len(document.tables) == 1
    table = document.tables[0]
    assert [[cell.text for cell in row.cells] for row in table.rows] == [
        ["Metric", "Value"],
        ["Revenue", "120"],
        ["Churn", "0.03"],
    ]
    assert "Appendix" in [p.text for p in document.paragraphs]


@pytest.mark.parametrize(
    ("payload", "message"),
    [
        (
            {"format": "docx", "blocks": [{"type": "heading", "level": 0, "text": "x"}]},
            "level",
        ),
        (
            {
                "format": "docx",
                "blocks": [
                    {"type": "table", "headers": ["a", "b"], "rows": [[1]]}
                ],
            },
            "same width",
        ),
        (
            {"format": "docx", "blocks": [{"type": "video", "url": "https://example.test"}]},
            "unsupported",
        ),
        (
            {"format": "docx", "blocks": [{"type": "bullet_list", "items": "not-a-list"}]},
            "array",
        ),
    ],
)
def test_render_docx_rejects_invalid_blocks(payload: object, message: str) -> None:
    with pytest.raises(OfficeDocumentError, match=message):
        render_office_document(payload)  # type: ignore[arg-type]


def test_render_xlsx_supports_multiple_sheets_and_neutralizes_formula_strings() -> None:
    rendered = render_office_document(
        {
            "format": "xlsx",
            "title": "Metrics workbook",
            "author": "Inkspan Agent",
            "sheets": [
                {
                    "name": "Summary",
                    "header_row": True,
                    "freeze_panes": "A2",
                    "auto_filter": True,
                    "rows": [
                        ["Metric", "Value", "Potential formula"],
                        ["Revenue", 120, "=1+1"],
                        ["Enabled", True, "@SUM(A1:A2)"],
                    ],
                },
                {"name": "Notes", "rows": [["Text"], [None], ["Done"]]},
            ],
        }
    )

    assert rendered.format == "xlsx"
    assert rendered.extension == ".xlsx"
    assert rendered.content_type.endswith("spreadsheetml.sheet")

    workbook = load_workbook(BytesIO(rendered.data), data_only=False)
    assert workbook.properties.title == "Metrics workbook"
    assert workbook.properties.creator == "Inkspan Agent"
    assert workbook.sheetnames == ["Summary", "Notes"]
    summary = workbook["Summary"]
    assert summary.freeze_panes == "A2"
    assert summary.auto_filter.ref == "A1:C3"
    assert summary["A1"].font.bold is True
    assert summary["B2"].value == 120
    assert summary["C2"].value == "=1+1"
    assert summary["C2"].data_type == "s"
    assert summary["C3"].value == "@SUM(A1:A2)"
    assert summary.column_dimensions["A"].width >= len("Revenue") + 2
    assert workbook["Notes"]["A2"].value is None


@pytest.mark.parametrize(
    ("payload", "message"),
    [
        ({"format": "xlsx", "sheets": []}, "at least one"),
        (
            {
                "format": "xlsx",
                "sheets": [
                    {"name": "Summary", "rows": []},
                    {"name": "summary", "rows": []},
                ],
            },
            "unique",
        ),
        (
            {"format": "xlsx", "sheets": [{"name": "bad/name", "rows": []}]},
            "invalid",
        ),
        (
            {
                "format": "xlsx",
                "sheets": [{"name": "Sheet", "rows": [[{"nested": "object"}]]}],
            },
            "scalar",
        ),
        (
            {
                "format": "xlsx",
                "sheets": [{"name": "Sheet", "rows": [], "header_row": "yes"}],
            },
            "boolean",
        ),
    ],
)
def test_render_xlsx_rejects_unsafe_or_ambiguous_workbooks(
    payload: object, message: str
) -> None:
    with pytest.raises(OfficeDocumentError, match=message):
        render_office_document(payload)  # type: ignore[arg-type]


def test_render_pptx_supports_title_and_content_slides() -> None:
    rendered = render_office_document(
        {
            "format": "pptx",
            "title": "Launch plan",
            "author": "Inkspan Agent",
            "slides": [
                {"title": "Launch plan", "subtitle": "August 2026"},
                {
                    "title": "Priorities",
                    "bullets": [
                        "Stabilize the API",
                        {"text": "Document the host contract", "level": 1},
                    ],
                },
                {"title": "Decision", "bullets": []},
            ],
        }
    )

    assert rendered.format == "pptx"
    assert rendered.extension == ".pptx"
    assert rendered.content_type.endswith("presentationml.presentation")

    presentation = Presentation(BytesIO(rendered.data))
    assert presentation.core_properties.title == "Launch plan"
    assert presentation.core_properties.author == "Inkspan Agent"
    assert len(presentation.slides) == 3
    assert presentation.slides[0].shapes.title.text == "Launch plan"
    assert presentation.slides[0].placeholders[1].text == "August 2026"
    body = presentation.slides[1].placeholders[1].text_frame
    assert [paragraph.text for paragraph in body.paragraphs] == [
        "Stabilize the API",
        "Document the host contract",
    ]
    assert [paragraph.level for paragraph in body.paragraphs] == [0, 1]
    assert presentation.slides[2].shapes.title.text == "Decision"


@pytest.mark.parametrize(
    ("payload", "message"),
    [
        ({"format": "pptx", "slides": []}, "at least one"),
        (
            {"format": "pptx", "slides": [{"title": "x", "subtitle": "s", "bullets": ["b"]}]},
            "subtitle",
        ),
        (
            {
                "format": "pptx",
                "slides": [{"title": "x", "bullets": [{"text": "b", "level": 9}]}],
            },
            "level",
        ),
        (
            {"format": "pptx", "slides": [{"title": "x", "bullets": [42]}]},
            "string or object",
        ),
    ],
)
def test_render_pptx_rejects_invalid_slide_shapes(payload: object, message: str) -> None:
    with pytest.raises(OfficeDocumentError, match=message):
        render_office_document(payload)  # type: ignore[arg-type]


@pytest.mark.parametrize(
    ("payload", "message"),
    [
        ([], "object"),
        ({}, "format"),
        ({"format": "pdf"}, "unsupported"),
        ({"format": 3}, "string"),
        ({"format": "docx", "title": 1, "blocks": []}, "title"),
    ],
)
def test_render_rejects_invalid_top_level_contract(payload: object, message: str) -> None:
    with pytest.raises(OfficeDocumentError, match=message):
        render_office_document(payload)  # type: ignore[arg-type]


def test_write_office_document_checks_extension_and_overwrite(tmp_path: Path) -> None:
    payload = {"format": "docx", "blocks": [{"type": "paragraph", "text": "hello"}]}
    output = tmp_path / "nested" / "document.docx"
    output.parent.mkdir(parents=True)
    unrelated_temporary = output.with_name(f".{output.name}.tmp")
    unrelated_temporary.write_text("keep", encoding="utf-8")

    written = write_office_document(payload, output)
    assert written == output
    assert output.read_bytes().startswith(b"PK")
    assert unrelated_temporary.read_text(encoding="utf-8") == "keep"

    with pytest.raises(FileExistsError):
        write_office_document(payload, output)

    write_office_document(payload, output, overwrite=True)

    with pytest.raises(OfficeDocumentError, match="extension"):
        write_office_document(payload, tmp_path / "document.xlsx")


def test_render_docx_supports_headerless_tables() -> None:
    rendered = render_office_document(
        {
            "format": "docx",
            "blocks": [{"type": "table", "rows": [["a", 1], ["b", 2]]}],
        }
    )
    document = Document(BytesIO(rendered.data))
    assert [[cell.text for cell in row.cells] for row in document.tables[0].rows] == [
        ["a", "1"],
        ["b", "2"],
    ]


@pytest.mark.parametrize(
    ("block", "message"),
    [
        ({"type": "table", "rows": []}, "at least one column"),
        ({"type": "table", "rows": [[1, 2], [3]]}, "same width"),
        ({"type": "heading", "level": True, "text": "x"}, "integer"),
        ({"type": "heading", "level": 1, "text": "   "}, "must not be empty"),
    ],
)
def test_render_docx_covers_remaining_validation_edges(
    block: dict[str, object], message: str
) -> None:
    with pytest.raises(OfficeDocumentError, match=message):
        render_office_document({"format": "docx", "blocks": [block]})


@pytest.mark.parametrize(
    ("payload", "message"),
    [
        ({"format": "docx"}, "blocks is required"),
        ({"format": "docx", "blocks": [], "bogus": 1}, "unexpected field"),
        (
            {"format": "docx", "blocks": [{"type": "paragraph", "text": "x", "bogus": 1}]},
            "unexpected field",
        ),
        (
            {"format": "xlsx", "sheets": [{"name": "Data"}]},
            "rows is required",
        ),
        (
            {"format": "xlsx", "sheets": [{"name": "Data", "rows": [], "bogus": 1}]},
            "unexpected field",
        ),
        (
            {"format": "pptx", "slides": [{"title": "x", "bogus": 1}]},
            "unexpected field",
        ),
        (
            {
                "format": "pptx",
                "slides": [{"title": "x", "bullets": [{"text": "y", "bogus": 1}]}],
            },
            "unexpected field",
        ),
    ],
)
def test_renderer_enforces_the_machine_readable_contract(
    payload: dict[str, object], message: str
) -> None:
    with pytest.raises(OfficeDocumentError, match=message):
        render_office_document(payload)


@pytest.mark.parametrize(
    ("payload", "message"),
    [
        ({"format": "DOCX", "blocks": []}, "unsupported format"),
        ({"format": "docx", "blocks": [], "title": None}, "title must be a string"),
        (
            {"format": "xlsx", "sheets": [{"name": "Data", "rows": [], "freeze_panes": None}]},
            "freeze_panes must be a string",
        ),
        (
            {"format": "pptx", "slides": [{"title": "Title", "subtitle": None}]},
            "subtitle must be a string",
        ),
        (
            {
                "format": "pptx",
                "slides": [{"title": "Title", "subtitle": "Sub", "bullets": []}],
            },
            "subtitle cannot be combined",
        ),
        (
            {"format": "xlsx", "sheets": [{"name": "Data", "rows": [[float("nan")]]}]},
            "finite",
        ),
        (
            {"format": "docx", "blocks": [], 1: "not-json"},
            "keys must be strings",
        ),
    ],
)
def test_renderer_rejects_values_outside_strict_json_schema(
    payload: dict[object, object], message: str
) -> None:
    with pytest.raises(OfficeDocumentError, match=message):
        render_office_document(payload)  # type: ignore[arg-type]
